"""PPTX ``DocumentParser`` implementation.

Walks slides via ``python-pptx``, converting each slide to Markdown:
title → ``#`` heading, text frames → paragraphs, tables → HTML→Markdown,
charts → Markdown tables, pictures → ``![](pptx-image-N)`` synthetic
markdown image refs. Speaker notes are appended as ``### Notes:``.

Captioning is not done here — see :class:`ImageBlock` for the
parser→caption contract.

Output is one :class:`TextBlock` per slide (1-indexed ``page_number``)
plus one :class:`ImageBlock` per slide picture. ``page_number`` on each
``ImageBlock`` is the slide number it came from.

Implementation derived from the legacy ``PPTXConverter`` (which itself
mirrored the MarkItDown PPTX converter).
"""

from __future__ import annotations

import asyncio
import html
import logging
from io import BytesIO
from typing import Any

from ...models.document import Document, DocumentType, ImageBlock, ProcessedDocument, TextBlock
from ..image_preprocessor import ensure_png_compatible_mode, pil_to_png_bytes
from .document_parser import DocumentParser

logger = logging.getLogger(__name__)


def _image_ref(index: int) -> str:
    """Synthetic markdown image ref used as a placeholder for slide pictures."""
    return f"![](pptx-image-{index})"


class PptxParser(DocumentParser):
    """Parse PPTX into one TextBlock per slide plus one ImageBlock per picture."""

    def supported_types(self) -> list[str]:
        return [DocumentType.PPTX.value]

    async def parse(self, document: Document) -> ProcessedDocument:
        if not document.raw_bytes:
            return ProcessedDocument(
                document_id=document.id,
                metadata=dict(document.metadata),
            )

        async with document.as_temporary_file() as path:
            slides, images = await asyncio.to_thread(self._convert, str(path))

        text_blocks = [TextBlock(text=text, page_number=slide_num) for slide_num, text in slides]
        page_count = slides[-1][0] if slides else 0
        return ProcessedDocument(
            document_id=document.id,
            text_blocks=text_blocks,
            images=images,
            metadata=dict(document.metadata),
            page_count=page_count,
        )

    # ----- conversion -----

    def _convert(self, path: str) -> tuple[list[tuple[int, str]], list[ImageBlock]]:
        try:
            import pptx
            from PIL import Image
        except ImportError:
            logger.warning("python-pptx or PIL not available; cannot parse PPTX")
            return "", []

        try:
            presentation = pptx.Presentation(path)
        except Exception as exc:
            logger.warning("Failed to open PPTX: %s", exc)
            return "", []

        slides: list[tuple[int, str]] = []
        images: list[ImageBlock] = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            md = ""
            title = slide.shapes.title

            for shape in slide.shapes:
                if self._is_picture(shape):
                    try:
                        with Image.open(BytesIO(shape.image.blob)) as im:
                            im = ensure_png_compatible_mode(im)
                            png_bytes = pil_to_png_bytes(im)
                        ref = _image_ref(len(images))
                        images.append(
                            ImageBlock(
                                image_bytes=png_bytes,
                                page_number=slide_num,
                                mime_type="image/png",
                                metadata={"markdown_ref": ref},
                            )
                        )
                        md += ref
                    except Exception as exc:
                        logger.warning("Skipping unreadable PPTX picture: %s", exc)
                elif self._is_table(shape):
                    md += "\n" + self._table_to_markdown(shape.table) + "\n"
                elif getattr(shape, "has_chart", False):
                    md += self._chart_to_markdown(shape.chart)
                elif getattr(shape, "has_text_frame", False):
                    if shape == title:
                        md += "# " + shape.text.lstrip() + "\n"
                    else:
                        md += shape.text + "\n"

            md = md.strip()
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md += "\n\n### Notes:\n" + notes_frame.text
                md = md.strip()

            if md:
                slides.append((slide_num, md))

        return slides, images

    @staticmethod
    def _is_picture(shape: Any) -> bool:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"):
                return True
        except NotImplementedError:
            logger.debug("Encountered an unimplemented shape type")
        except Exception:
            return False
        return False

    @staticmethod
    def _is_table(shape: Any) -> bool:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            return shape.shape_type == MSO_SHAPE_TYPE.TABLE
        except NotImplementedError:
            logger.debug("Encountered an unimplemented shape type")
            return False
        except Exception:
            return False

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        from html_to_markdown import convert

        html_rows = ["<html><body><table>"]
        first_row = True
        for row in table.rows:
            html_rows.append("<tr>")
            for cell in row.cells:
                tag = "th" if first_row else "td"
                html_rows.append(f"<{tag}>{html.escape(cell.text)}</{tag}>")
            html_rows.append("</tr>")
            first_row = False
        html_rows.append("</table></body></html>")
        return convert("".join(html_rows)).strip()

    @staticmethod
    def _chart_to_markdown(chart: Any) -> str:
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data: list[list[str]] = [["Category"] + series_names]
            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)
            rows = ["| " + " | ".join(map(str, r)) + " |" for r in data]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([rows[0], separator] + rows[1:])
        except ValueError as exc:
            if "unsupported plot type" in str(exc):
                return "\n\n[unsupported chart]\n\n"
            return "\n\n[unsupported chart]\n\n"
        except Exception:
            return "\n\n[unsupported chart]\n\n"
