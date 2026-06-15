import html
import re
from io import BytesIO

import pptx
from html_to_markdown import convert
from langchain_core.documents.base import Document
from PIL import Image
from utils.logger import get_logger

from .base import BaseLoader

logger = get_logger()


class PPTXConverter:
    """Implementation based on PPTX converter in MarkItDown library.

    https://github.com/microsoft/markitdown/blob/main/packages/markitdown/src/markitdown/converters/_pptx_converter.py
    """

    def __init__(
        self,
        image_placeholder=r"<image>",
        page_separator: str = "[PAGE_SEP]",
        max_slides: int = 2000,
        max_images: int = 2000,
    ):
        self.image_placeholder = image_placeholder
        self.page_separator = page_separator
        # Parser-bomb caps: bound slides processed and images decoded into memory.
        self.max_slides = max_slides
        self.max_images = max_images

    def convert(self, local_path):
        md_content = ""
        presentation = pptx.Presentation(local_path)
        slide_num = 0
        images_list = []
        for slide in presentation.slides:
            slide_num += 1
            if slide_num > self.max_slides:
                logger.warning("Capping slide processing", cap=self.max_slides)
                break

            title = slide.shapes.title
            for shape in slide.shapes:
                if self._is_picture(shape):
                    if len(images_list) < self.max_images:
                        images_list.append(Image.open(BytesIO(shape.image.blob)))
                        md_content += self.image_placeholder

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += "\n" + convert(html_table).strip() + "\n"

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

            md_content += f"\n[PAGE_{slide_num}]\n"

        return md_content, images_list

    def _is_picture(self, shape):
        try:
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                return True
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "image"):
                    return True
        except NotImplementedError:
            # https://python-pptx.readthedocs.io/en/latest/_modules/pptx/shapes/autoshape.html
            # Not all shape types are implemented in python-pptx
            logger.warning("Encountered an unimplemented shape type.")

        return False

    def _is_table(self, shape):
        try:
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
                return True
        except NotImplementedError:
            # # https://python-pptx.readthedocs.io/en/latest/_modules/pptx/shapes/autoshape.html
            # Not all shape types are implemented in python-pptx
            logger.warning("Encountered an unimplemented shape type.")
        return False

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"


class PPTXLoader(BaseLoader):
    def __init__(self, **kwargs) -> None:
        super().__init__(**kwargs)
        self.image_placeholder = r"<image>"
        max_entries = int(self.config.loader.get("max_archive_entries", 2000))
        self.converter = PPTXConverter(
            image_placeholder=self.image_placeholder,
            page_separator=self.page_sep,
            max_slides=max_entries,
            max_images=max_entries,
        )

    async def aload_document(self, file_path, metadata=None, save_markdown=False):
        md_content, imgs = self.converter.convert(local_path=file_path)

        if self.image_captioning:
            images_captions = await self.caption_images(imgs, desc="Generating captions")

            for caption in images_captions:
                md_content = re.sub(
                    self.image_placeholder,
                    caption.replace("\\", "/"),
                    md_content,
                    count=1,
                )
        else:
            logger.info("Image captioning disabled. Ignoring images.")
            # Remove image placeholders when captioning is disabled
            md_content = md_content.replace(self.image_placeholder, "")

        doc = Document(page_content=md_content, metadata=metadata)
        if save_markdown:
            self.save_content(md_content, str(file_path))
        return doc
